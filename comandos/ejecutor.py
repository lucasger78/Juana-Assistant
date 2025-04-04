# import os
# import subprocess
# import time
# from django.http import JsonResponse

# def ejecutar(accion: str) -> str:
#     accion = accion.lower()
    
#     try:
#         if "chrome" in accion:
#             # Método 100% funcional en Windows
#             os.system('start chrome --new-window')
#             return "✅ Chrome abierto"
            
#         elif "spotify" in accion:
#             try:
#                 # Método 100% funcional
#                 spotify_path = os.path.join(os.environ['APPDATA'], 'Spotify', 'Spotify.exe')
#                 subprocess.Popen([spotify_path], shell=True)
                
#                 # Verificación adicional
#                 time.sleep(2)  # Espera 2 segundos
#                 procesos = subprocess.check_output('tasklist', shell=True).decode()
#                 if "Spotify.exe" in procesos:
#                     return "🎵 Spotify abierto correctamente"
#                 else:
#                     raise Exception("No se detectó en memoria")
                    
#             except Exception as e:
#                 # Fallback 1: Protocolo spotify://
#                 subprocess.Popen(['start', 'spotify://'], shell=True)
                
#                 # Fallback 2: Versión web
#                 subprocess.Popen(['start', 'https://open.spotify.com'], shell=True)
#                 return "🌐 Spotify Web abierto"
            
#         elif "calculadora" in accion:
#             os.system('calc')
#             return "🧮 Calculadora abierta"
            
#         elif "bloc de notas" in accion:
#             os.system('notepad')
#             return "📝 Bloc de notas listo"
        
       
            
#     except Exception as e:
#         return f"❌ Error: {str(e)}"
    
#     return "⚠️ Comando no reconocido"
from django.shortcuts import render
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt

import os
import time
import speech_recognition as sr
import pyautogui
import urllib.parse
import webbrowser
import re

from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.common.exceptions import TimeoutException, NoSuchElementException
from selenium.webdriver.chrome.options import Options
from reportlab.pdfgen import canvas
from openpyxl import Workbook
from pptx import Presentation
from docx import Document
from pathlib import Path
from pathlib import Path
from reportlab.pdfgen import canvas
from openpyxl import Workbook
from pptx import Presentation
from docx import Document

# Configuración
WHATSAPP_WAIT_TIME = 10  # Tiempo máximo de espera para que cargue WhatsApp Web
CHROME_PROFILE_PATH = "C:\\chrome_selenium"  # Ruta del perfil de usuario de Chrome
REMOTE_DEBUGGING_PORT = "9222"

def home(request):
    """Vista para la página principal"""
    return render(request, 'index.html')

@csrf_exempt
def escuchar(request):
    """Procesa comandos de voz"""
    if request.method == 'POST':
        try:
            # Reconocimiento de voz
            r = sr.Recognizer()
            with sr.Microphone() as source:
                audio = r.listen(source, timeout=5)
            texto = r.recognize_google(audio, language="es-ES").lower().strip()
            print(f"🧪 ACCIÓN DETECTADA (escuchar): '{texto}'")
            print(f"Comando detectado: {texto}")  # Debugging     
            
            # Ejecutar comando
            resultado = ejecutar_comando(texto)
            return JsonResponse({"status": resultado, "comando": texto})
            
        except sr.WaitTimeoutError:
            return JsonResponse({"error": "No se detectó voz"}, status=400)
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "Método no permitido"}, status=405)

def obtener_ruta_escritorio():
    """Devuelve la ruta del escritorio del usuario"""
    return os.path.join(os.path.expanduser("~"), "Desktop")

def crear_archivo(accion: str) -> str:
    """Crea un archivo según el comando recibido (ej: 'crear archivo de excel notas')"""
    accion = accion.lower().strip()

    # Detectar tipo de archivo
    if "pdf" in accion:
        extension = ".pdf"
    elif "word" in accion:
        extension = ".docx"
    elif "excel" in accion:
        extension = ".xlsx"
    elif "powerpoint" in accion or "presentación" in accion:
        extension = ".pptx"
    elif "txt" in accion or "texto" in accion:
        extension = ".txt"
    else:
        return "⚠️ Tipo de archivo no reconocido. Usa: PDF, Word, Excel, PowerPoint o TXT."

    # Extraer nombre del archivo
    partes = accion.split()
    if "de" in partes:
        index = partes.index("de")
        nombre = "_".join(partes[index + 2:])  # Tomar el nombre después de "de <tipo> ..."
    else:
        return "⚠️ No se detectó un nombre válido para el archivo."

    if not nombre:
        return "⚠️ Especifica un nombre para el archivo."

    # Generar ruta con extensión correcta
    ruta_base = obtener_ruta_escritorio()
    ruta = os.path.join(ruta_base, nombre + extension)

    # Crear archivo según su tipo
    try:
        if extension == ".txt":
            with open(ruta, 'w') as archivo:
                archivo.write("")
            return f"📄 Archivo TXT '{nombre}' creado en {ruta}"

        elif extension == ".pdf":
            c = canvas.Canvas(ruta)
            c.drawString(100, 750, "Nuevo archivo PDF")
            c.save()
            return f"📄 Archivo PDF '{nombre}' creado en {ruta}"

        elif extension == ".docx":
            doc = Document()
            doc.add_paragraph("Nuevo archivo Word")
            doc.save(ruta)
            return f"📄 Archivo Word '{nombre}' creado en {ruta}"

        elif extension == ".xlsx":
            wb = Workbook()
            ws = wb.active
            ws["A1"] = "Nuevo archivo Excel"
            wb.save(ruta)
            return f"📊 Archivo Excel '{nombre}' creado en {ruta}"

        elif extension == ".pptx":
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            textbox = slide.shapes.add_textbox(100, 100, 500, 100)
            textbox.text = "Nuevo archivo PowerPoint"
            prs.save(ruta)
            return f"📽️ Archivo PowerPoint '{nombre}' creado en {ruta}"

    except Exception as e:
        return f"❌ Error al crear el archivo: {str(e)}"

    return "⚠️ Algo salió mal."



import webbrowser

def ejecutar_comando(accion: str) -> str:
    accion_original = accion
    accion = accion.strip().lower()

    print(f"🧪 ACCIÓN ORIGINAL: '{accion_original}'")
    print(f"🔍 ACCIÓN PROCESADA: '{accion}'")
    print(f"🧬 Longitud: {len(accion)}")
    print(f"🔠 Caracteres (código ASCII): {[ord(c) for c in accion]}")

    # Comandos directos (simples)
    comandos = {
        "whatsapp": "https://web.whatsapp.com",
        "youtube": "https://www.youtube.com",
        "gmail": "https://mail.google.com",
        "flow": "https://portal.app.flow.com.ar/inicio",
        "disney": "https://www.disneyplus.com/es-419/home",
        "max": "https://play.max.com/",
        "wikipedia": "https://www.wikipedia.org",
        "linkedin": "https://www.linkedin.com",
        "facebook": "https://www.faceook.com",
        "intagram": "https://www.instagram.com",
        "pinterest": "https://ar.pinterest.com/",
        "x": "https://x.com/?lang=es",
        "reddit": "https://www.reddit.com/r/argentina/",
        "netflix": "https://www.netflix.com",
        "noticias": "https://news.google.com",
        "calendar": "https://calendar.google.com/calendar/u/0/r/eventedit"
    }

    for clave, url in comandos.items():
        print(f"🧩 Comparando con: '{clave}'")
        if clave in accion:
            print(f"🎯 COINCIDENCIA: '{clave}' → {url}")
            webbrowser.open(url)
            return f"✅ {clave.capitalize()} abierto"

    # Comando: Buscar en Google
    if "buscar en google" in accion:
        query = accion.split("buscar en google")[-1].strip()
        if query:
            webbrowser.open(f"https://www.google.com/search?q={query}")
            return f"🔎 Buscando '{query}' en Google"
        return "🧐 ¿Qué querés buscar?"

    # Comando: Reproducir video de [X] en YouTube
    if "reproducir video de" in accion and "en youtube" in accion:
        busqueda = accion.split("reproducir video de")[-1].split("en youtube")[0].strip()
        if busqueda:
            webbrowser.open(f"https://www.youtube.com/results?search_query={busqueda}")
            return f"📺 Buscando '{busqueda}' en YouTube"
        return "❓ ¿Qué querés ver en YouTube?"

    # Comando: Enviar correo a [nombre]
    if "enviar correo a" in accion:
        destinatario = accion.split("enviar correo a")[-1].strip()
        if destinatario:
            webbrowser.open(f"mailto:{destinatario}")
            return f"📧 Redactando correo a {destinatario}"
        return "📭 ¿A quién querés escribirle?"

    print("🚫 No se encontró ningún comando válido.")
    return f"⚠️ Comando no reconocido: {accion}"





def abrir_en_chrome(url: str):
    """Abre una URL en Chrome con la sesión existente"""
    try:
        options = Options()
        options.add_argument(f"--remote-debugging-port={REMOTE_DEBUGGING_PORT}")
        driver = webdriver.Chrome(options=options)
        driver.get(url)
    except Exception as e:
        print(f"❌ Error al abrir {url}: {str(e)}")

#############################################################################################
#############################################################################################
#############################################################################################

def enviar_mensaje_whatsapp(accion: str) -> str:
    """Envía un mensaje por WhatsApp Web reutilizando la sesión abierta"""
    try:
        partes = accion.split("a") if "a " in accion else accion.split("para")
        nombre_contacto = partes[-1].split(":")[0].strip()
        mensaje = ":".join(partes[-1].split(":")[1:]).strip()
        
        options = Options()
        options.debugger_address = f"127.0.0.1:{REMOTE_DEBUGGING_PORT}"
        driver = webdriver.Chrome(options=options)
        driver.get("https://web.whatsapp.com")
        
        try:
            WebDriverWait(driver, WHATSAPP_WAIT_TIME).until(
                EC.presence_of_element_located((By.XPATH, '//div[@contenteditable="true"][@data-tab="3"]'))
            )
        except TimeoutException:
            return "⚠️ WhatsApp no cargó a tiempo. Verifica la conexión."

        search_box = driver.find_element(By.XPATH, '//div[@contenteditable="true"][@data-tab="3"]')
        search_box.click()
        search_box.clear()
        search_box.send_keys(nombre_contacto)
        time.sleep(2)

        try:
            driver.find_element(By.XPATH, f'//span[@title="{nombre_contacto}"]').click()
        except NoSuchElementException:
            return f"⚠️ No se encontró el contacto: {nombre_contacto}"

        try:
            input_box = WebDriverWait(driver, 5).until(
                EC.presence_of_element_located((By.XPATH, '//div[@contenteditable="true"][@data-tab="1"]'))
            )
            input_box.send_keys(mensaje)
            pyautogui.press('enter')
            return f"📤 Mensaje enviado a {nombre_contacto.capitalize()}"
        except TimeoutException:
            return "⚠️ No se pudo encontrar el campo de mensaje."

    except Exception as e:
        return f"❌ Error WhatsApp: {str(e)}"
    
    
def obtener_ruta_escritorio():
    """Devuelve la ruta del escritorio del usuario"""
    return os.path.join(os.path.expanduser("~"), "Desktop")

def crear_carpeta(nombre: str) -> str:
    """Crea una carpeta en el escritorio"""
    ruta_base = obtener_ruta_escritorio()
    ruta = os.path.join(ruta_base, nombre)
    os.makedirs(ruta, exist_ok=True)
    return f"📁 Carpeta '{nombre}' creada en {ruta}"

# def crear_archivo(nombre: str) -> str:
#     """Crea un archivo en el escritorio"""
#     ruta_base = obtener_ruta_escritorio()
#     ruta = os.path.join(ruta_base, nombre)
#     with open(ruta, 'w') as archivo:
#         archivo.write("")
#     return f"📄 Archivo '{nombre}' creado en {ruta}"
#   SOLO PARA TXT


    
    
